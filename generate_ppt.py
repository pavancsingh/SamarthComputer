import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_PPTX = r"d:\Samarthcomputers\docs\project-presentation\Samarth_Computers_BCA_Presentation.pptx"
SCREENSHOT_DIR = r"d:\Samarthcomputers\docs\project-presentation\screenshots"

# Color Palette
COLOR_PRIMARY = RGBColor(198, 40, 40)    # Samarth Red #C62828
COLOR_NAVY = RGBColor(15, 23, 42)       # Navy Slate #0F172A
COLOR_DARK_TEXT = RGBColor(30, 41, 59)   # Dark Gray #1E293B
COLOR_MUTED = RGBColor(100, 116, 139)   # Muted Gray #64748B
COLOR_BG_CARD = RGBColor(241, 245, 249)  # Light Slate #F1F5F9
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_EMERALD = RGBColor(16, 185, 129)   # Emerald #10B981
COLOR_BORDER = RGBColor(226, 232, 240)  # Border #E2E8F0

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="BCA PROJECT PRESENTATION"):
    # Header bar background / title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_PRIMARY

    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_NAVY
    p_title.space_before = Pt(4)

def add_card(slide, left, top, width, height, title="", body_items=[], bg_color=COLOR_BG_CARD, border_color=COLOR_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    if title:
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        p_title.space_after = Pt(10)
        first_bullet = True
    else:
        first_bullet = False

    for item in body_items:
        if first_bullet and not title:
            p = tf.paragraphs[0]
            first_bullet = False
        else:
            p = tf.add_paragraph()
        
        p.text = "• " + item if isinstance(item, str) else item.get('text', '')
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.space_after = Pt(6)

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text

def try_add_image(slide, image_filename, left, top, width, height):
    img_path = os.path.join(SCREENSHOT_DIR, image_filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width, height)
        return True
    return False

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # SLIDE 1: TITLE
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_NAVY)

    # Accent bar
    accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.2), Inches(3.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_PRIMARY
    accent.line.fill.background()

    txBox = slide1.shapes.add_textbox(Inches(1.3), Inches(2.1), Inches(11.0), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "BCA FINAL YEAR ACADEMIC PROJECT PRESENTATION"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_EMERALD

    p1 = tf.add_paragraph()
    p1.text = "Samarth Computers Khandala"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_before = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "Modern Computer Education & Digital Services Platform"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    p2.space_before = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "MKCL Authorized Learning Center (ALC 13210399 / 13210273) | Full-Stack Web Application"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_before = Pt(16)

    add_speaker_notes(slide1, "Good morning respected examiners and faculty members. Today I am presenting my BCA final year project: Samarth Computers Khandala, a full-stack digital platform for computer training and government service assistance.")

    # -------------------------------------------------------------
    # SLIDE 2: PROJECT OVERVIEW
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_WHITE)
    add_header(slide2, "Project Overview & Background")

    add_card(slide2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             title="What is Samarth Computers?",
             body_items=[
                 "Established MKCL Authorized Learning Center in Khandala, Satara.",
                 "Offers IT courses: MS-CIT, Tally Prime GST, Typing (GCC-TBC), Advanced Excel.",
                 "Operates an official CSC Digital Seva Kendra for citizen services."
             ])

    add_card(slide2, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0),
             title="Purpose of Digital Platform",
             body_items=[
                 "Bridge student aspirations with 24/7 web course discovery.",
                 "Eliminate document confusion for government service applicants.",
                 "Provide automated student certificate verification."
             ])

    add_card(slide2, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0),
             title="System Architecture Scope",
             body_items=[
                 "High-converting public marketing website.",
                 "Comprehensive Admin Operations Dashboard.",
                 "Real-time data synchronization backed by Supabase PostgreSQL."
             ])

    add_speaker_notes(slide2, "Samarth Computers is a trusted IT institute in Khandala. This project delivers a production-grade web platform combining a public portal for course/CSC discovery with a full Admin Dashboard.")

    # -------------------------------------------------------------
    # SLIDE 3: PROBLEM STATEMENT
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_WHITE)
    add_header(slide3, "Problem Statement: Challenges in Traditional Management")

    add_card(slide3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3),
             title="Manual Content & Pamphlet Dependency",
             body_items=[
                 "Course syllabus, fee details, and batch timings relied on printed flyers and verbal explanation.",
                 "No single digital destination for prospective students to review course curriculum."
             ])

    add_card(slide3, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.3),
             title="Difficult Service & Fee Updates",
             body_items=[
                 "Updating batch timetables or fees required physical poster re-printing.",
                 "Static websites required developer intervention for basic text modifications."
             ])

    add_card(slide3, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.4),
             title="Scattered Student & Customer Data",
             body_items=[
                 "Student admission inquiries were written manually in physical registers.",
                 "High risk of lost phone numbers, unorganized leads, and delayed follow-ups."
             ])

    add_card(slide3, Inches(6.8), Inches(4.4), Inches(5.7), Inches(2.4),
             title="Limited Citizen Guidance & Accessibility",
             body_items=[
                 "Citizens frequently arrived at CSC center without required application documents.",
                 "No online mechanism for employers to verify course completion certificates."
             ])

    add_speaker_notes(slide3, "Before this application, institute inquiries were logged in paper registers, document checklists were explained verbally, and content updates required manual manual effort.")

    # -------------------------------------------------------------
    # SLIDE 4: OBJECTIVES
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_WHITE)
    add_header(slide4, "Project Objectives")

    objectives = [
        ("1. Digitalize Institute Information", "Provide 24/7 web access to MKCL course syllabi, fee schedules, faculty profiles, and campus facilities."),
        ("2. Streamline CSC & Govt Services", "Eliminate citizen confusion with interactive document checklists and step-by-step application instructions."),
        ("3. Centralize Administration", "Build a secure Admin Dashboard featuring rate-limited auth, lead inbox management, and 19 CRUD sub-panels."),
        ("4. Enable Certificate Verification", "Implement instant online verification of student completion certificates using registration numbers."),
        ("5. Responsive & Bilingual Experience", "Deliver zero-reload dynamic language switching between Marathi and English with mobile-first layouts.")
    ]

    top_pos = 1.8
    for title, desc in objectives:
        add_card(slide4, Inches(0.8), Inches(top_pos), Inches(11.7), Inches(0.95),
                 title=title, body_items=[desc])
        top_pos += 1.05

    add_speaker_notes(slide4, "Our core objectives were to digitize institute communications, streamline CSC document guidance, centralize administration, provide certificate verification, and deliver a bilingual Marathi/English user experience.")

    # -------------------------------------------------------------
    # SLIDE 5: PROPOSED SYSTEM
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLOR_WHITE)
    add_header(slide5, "Proposed System & Operational Flow")

    # Flow visual boxes
    add_card(slide5, Inches(0.8), Inches(2.2), Inches(3.6), Inches(4.5),
             title="Public Visitor Workflow",
             body_items=[
                 "Explore MKCL Courses & Syllabus Modules",
                 "Review CSC Document Checklists",
                 "Check Live Batch Timetables",
                 "Verify Student Completion Certificates",
                 "Submit Admission & Service Inquiry Leads"
             ])

    add_card(slide5, Inches(4.8), Inches(2.2), Inches(3.6), Inches(4.5),
             title="React 19 App Shell",
             body_items=[
                 "Client-side SPA Router (currentView)",
                 "Bilingual Engine (Marathi / English)",
                 "Stitch Design System Tokens",
                 "AuthContext & Security Rate Limiter",
                 "sharedStore LocalStorage Observer Cache"
             ])

    add_card(slide5, Inches(8.8), Inches(2.2), Inches(3.7), Inches(4.5),
             title="Admin & Database Backend",
             body_items=[
                 "Rate-limited Admin Login Authentication",
                 "Real-Time Digital Leads Inbox",
                 "19 Modular CRUD Control Panels",
                 "Supabase PostgreSQL (10 RLS Tables)",
                 "Supabase Storage (samarth-media bucket)"
             ])

    add_speaker_notes(slide5, "The proposed system replaces manual logbooks with a unified React single-page application connected directly to Supabase cloud infrastructure for real-time lead tracking and content management.")

    # -------------------------------------------------------------
    # SLIDE 6: TECHNOLOGY STACK
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLOR_WHITE)
    add_header(slide6, "Technology Stack & Frameworks")

    tech_stack = [
        ("React 19", "Frontend UI Library", "Client-side rendering, component-driven UI architecture."),
        ("Vite 6", "Build Tool & Dev Server", "Ultra-fast HMR compilation, optimized production bundling."),
        ("Tailwind CSS 3.4", "Design System", "Stitch design tokens, responsive grid layouts, custom typography."),
        ("Supabase Cloud", "Backend Infrastructure", "PostgreSQL database, Auth, Storage bucket, instant APIs."),
        ("PostgreSQL", "Relational Database", "10 tables with Row Level Security (RLS) policies."),
        ("Cloudflare Pages", "Hosting & Edge CDN", "Global high-availability deployment with SSL & deep-link routing.")
    ]

    positions = [(0.8, 1.8), (4.8, 1.8), (8.8, 1.8), (0.8, 4.4), (4.8, 4.4), (8.8, 4.4)]
    for i, (title, tag, desc) in enumerate(tech_stack):
        left, top = positions[i]
        add_card(slide6, Inches(left), Inches(top), Inches(3.7), Inches(2.4),
                 title=f"{title} ({tag})", body_items=[desc])

    add_speaker_notes(slide6, "We selected React 19 and Vite 6 for high-speed client rendering, Tailwind CSS 3.4 for UI styling, Supabase PostgreSQL for backend persistence, and Cloudflare Pages for edge hosting.")

    # -------------------------------------------------------------
    # SLIDE 7: WEBSITE FEATURES
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLOR_WHITE)
    add_header(slide7, "Public Website Features")

    if not try_add_image(slide7, "01-home-hero.png", Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8)):
        add_card(slide7, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), title="Home Page Landing View", body_items=["Hero Banner, Accreditation Badges, Statistics Counter"])

    add_card(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8),
             title="Key User Modules",
             body_items=[
                 "Home Page: Hero banner, MKCL badges, statistics counter.",
                 "Course Catalog: Filterable MS-CIT, Tally, Typing, Excel cards.",
                 "CSC Services Desk: Identity & Revenue scheme guides.",
                 "Faculty Directory: Teacher credentials & experience.",
                 "Batch Timetable: Morning/Afternoon/Evening schedule slots.",
                 "Certificate Verification: Search by Registration Number.",
                 "Contact Page: Address, click-to-call, embedded map."
             ])

    add_speaker_notes(slide7, "The public website provides an intuitive interface for students and citizens to explore courses, check timetables, verify certificates, and locate the institute.")

    # -------------------------------------------------------------
    # SLIDE 8: COURSES & CSC SERVICES
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLOR_WHITE)
    add_header(slide8, "Courses & CSC Government Services Desk")

    if not try_add_image(slide8, "03-course-details-modal.png", Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8)):
        if not try_add_image(slide8, "04-csc-services-desk.png", Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8)):
            add_card(slide8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), title="Interactive Modals View", body_items=["Syllabus breakdown & Document Checklists"])

    add_card(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8),
             title="Interactive Discovery",
             body_items=[
                 "Categorized Course Exploration: MS-CIT, Tally Prime GST, Typing (GCC-TBC 30/40 wpm), Advanced Excel.",
                 "Course Detail Modal: Syllabus modules list, duration, fee estimates, and career opportunities.",
                 "CSC Identity Services: Aadhaar Card, PAN Card, Voter ID.",
                 "Revenue Services: 7/12 Utara, Domicile, Income Certificate.",
                 "Document Checklist Modal (DocChecklistModal): Clear step-by-step document guidance for local citizens."
             ])

    add_speaker_notes(slide8, "Students can open detailed syllabus drawers for any course, while citizens receive exact document checklists before visiting the center for government certificate applications.")

    # -------------------------------------------------------------
    # SLIDE 9: ADMIN DASHBOARD
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, COLOR_WHITE)
    add_header(slide9, "Centralized Admin Operations Dashboard")

    if not try_add_image(slide9, "12-admin-overview.png", Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8)):
        add_card(slide9, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), title="Admin Overview Screen", body_items=["KPI Cards, Quick Action Bar, Live Leads Feed"])

    add_card(slide9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8),
             title="Management Capabilities",
             body_items=[
                 "Overview Analytics: Live KPI cards for total leads, active courses, CSC services, and faculty count.",
                 "Quick Action Bar: One-click creation for courses, timetable slots, gallery photos, and news.",
                 "Real-Time Leads Inbox: Captures student inquiries with status tags (New Lead, In Process, Completed).",
                 "Live Database Persistence: All administrative edits sync directly with Supabase PostgreSQL."
             ])

    add_speaker_notes(slide9, "The Admin Dashboard gives institute managers total control over student lead tracking, course catalogs, timetable slots, and center announcements from a single screen.")

    # -------------------------------------------------------------
    # SLIDE 10: ADMIN SIDEBAR ARCHITECTURE
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, COLOR_WHITE)
    add_header(slide10, "Admin Sidebar Architecture & Accordion Redesign")

    if not try_add_image(slide10, "15-admin-sidebar-redesign.png", Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8)):
        add_card(slide10, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), title="Redesigned 3-Tier Sidebar", body_items=["Main Nav, Website Pages, Settings Accordion"])

    add_card(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8),
             title="3-Tier Structure",
             body_items=[
                 "Tier 1: MAIN NAVIGATION - Dashboard, Inbox Leads, Courses, CSC Services, Faculty, News, Timetable, Photos.",
                 "Tier 2: WEBSITE PAGES - Home Page, About Page, Contact & Call CTA configuration.",
                 "Tier 3: SETTINGS & CONFIG ACCORDION - Consolidated single collapsible menu housing 10 sub-panels:",
                 "  └ Navigation, Branding, Govt Certificates, Theme, Site Info, SEO, Social Links, Footer."
             ])

    add_speaker_notes(slide10, "To improve admin usability, we restructured the sidebar into three distinct categories and consolidated secondary site configuration controls into a clean expandable accordion.")

    # -------------------------------------------------------------
    # SLIDE 11: SUPABASE DATABASE ARCHITECTURE
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, COLOR_WHITE)
    add_header(slide11, "Supabase Database & Data Flow Architecture")

    add_card(slide11, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8),
             title="React 19 Frontend",
             body_items=[
                 "App Shell (App.jsx)",
                 "State View Router",
                 "AuthContext Rate Limiter",
                 "Stitch UI Components",
                 "Bilingual Prop State"
             ])

    add_card(slide11, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8),
             title="Repository Layer",
             body_items=[
                 "AdminRepository.js",
                 "CourseRepository.js",
                 "InquiryRepository.js",
                 "sharedStore (Cache)",
                 "StorageService.js"
             ])

    add_card(slide11, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8),
             title="Supabase Cloud",
             body_items=[
                 "PostgreSQL Engine",
                 "10 RLS Tables",
                 "Supabase Auth",
                 "samarth-media Bucket",
                 "Base64 Fallback Engine"
             ])

    add_speaker_notes(slide11, "The application follows a clean 3-tier data architecture linking the React client to repository abstractions, which query Supabase cloud database and storage services.")

    # -------------------------------------------------------------
    # SLIDE 12: DATABASE MODULES
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, COLOR_WHITE)
    add_header(slide12, "Database Modules: 10 Verified PostgreSQL Tables")

    tables = [
        ("inquiries", "Leads inbox storing student admissions & service inquiries."),
        ("courses", "Course catalog with JSONB syllabus modules & fee structure."),
        ("csc_services", "CSC identity services with JSONB document checklists."),
        ("govt_services", "Revenue services with JSONB application steps."),
        ("faculties", "Instructor profiles, designations, & specializations."),
        ("site_gallery", "Campus lab photos & event gallery records."),
        ("site_settings", "Central center configuration record (main_settings)."),
        ("batches", "Morning, afternoon, evening timetable slots."),
        ("news", "Institute announcements & examination news updates."),
        ("certificates", "Student completion records for online verification.")
    ]

    positions = [(0.8, 1.8), (6.8, 1.8), (0.8, 2.8), (6.8, 2.8), (0.8, 3.8), (6.8, 3.8), (0.8, 4.8), (6.8, 4.8), (0.8, 5.8), (6.8, 5.8)]
    for i, (tbl, desc) in enumerate(tables):
        left, top = positions[i]
        add_card(slide12, Inches(left), Inches(top), Inches(5.7), Inches(0.85),
                 title=f"table: public.{tbl}", body_items=[desc])

    add_speaker_notes(slide12, "The backend relies on 10 verified PostgreSQL tables in Supabase. We utilize JSONB columns for flexible syllabus modules and document checklists.")

    # -------------------------------------------------------------
    # SLIDE 13: SETTINGS & CONFIGURATION
    # -------------------------------------------------------------
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13, COLOR_WHITE)
    add_header(slide13, "Site Settings & Configuration Persistence")

    if not try_add_image(slide13, "16-admin-branding-settings.png", Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8)):
        add_card(slide13, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), title="Settings Control Panel", body_items=["Branding, Logo Upload, SEO, Social Links"])

    add_card(slide13, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8),
             title="Configurable Sub-Panels",
             body_items=[
                 "Navigation: Custom menu order & visibility.",
                 "Branding: Institute logo upload & hero banner background.",
                 "Govt Certificates: MKCL ALC code & CSC registration ID.",
                 "Theme: Custom primary & secondary color accents.",
                 "Site Information: Office hours, contact numbers, map URL.",
                 "SEO & Meta: Search meta title, description, and keywords.",
                 "Social Links: Facebook, Instagram, YouTube URLs.",
                 "Footer: Tagline, copyright notice, & quick links."
             ])

    add_speaker_notes(slide13, "All center parameters, logos, contact numbers, and SEO metadata are managed dynamically through the site settings accordion and saved to Supabase.")

    # -------------------------------------------------------------
    # SLIDE 14: SECURITY
    # -------------------------------------------------------------
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14, COLOR_WHITE)
    add_header(slide14, "Security & Data Protection Engine")

    add_card(slide14, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3),
             title="Admin Auth & Brute-Force Protection",
             body_items=[
                 "Allowed Emails List: Access restricted to authorized center emails (pawansingh3760@gmail.com).",
                 "Failed-Login Rate Limiter: Lockout timer triggers after 5 failed password attempts."
             ])

    add_card(slide14, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.3),
             title="PostgreSQL Row Level Security (RLS)",
             body_items=[
                 "RLS enabled across all 10 database tables.",
                 "Public SELECT for public catalog tables; Public INSERT strictly for lead inquiries."
             ])

    add_card(slide14, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.4),
             title="Session Management",
             body_items=[
                 "Encapsulated AuthContext provider.",
                 "Persistent encrypted session tokens stored in LocalStorage (samarth_admin_session)."
             ])

    add_card(slide14, Inches(6.8), Inches(4.4), Inches(5.7), Inches(2.4),
             title="Environment & Storage Security",
             body_items=[
                 "VITE_SUPABASE_URL & ANON_KEY environment isolation.",
                 "Protected storage bucket upload policies (samarth-media)."
             ])

    add_speaker_notes(slide14, "Security is enforced at multiple levels: rate-limited login, allowed email validation, LocalStorage session persistence, and PostgreSQL Row Level Security policies.")

    # -------------------------------------------------------------
    # SLIDE 15: RESPONSIVE & BILINGUAL EXPERIENCE
    # -------------------------------------------------------------
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15, COLOR_WHITE)
    add_header(slide15, "Responsive UI & Bilingual (Marathi / English) Engine")

    if not try_add_image(slide15, "17-mobile-responsive-ui.png", Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.8)):
        add_card(slide15, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.8), title="Mobile Responsive Viewport", body_items=["Smartphone Screen & Drawer Menu"])

    add_card(slide15, Inches(0.8), Inches(1.8), Inches(6.6), Inches(4.8),
             title="Accessibility & Globalization",
             body_items=[
                 "Dynamic Language Switching: Zero-reload toggle between Marathi (मराठी) and English across all UI components.",
                 "Dual-Language Data Model: Backend schemas support dual text fields (title_mr / title_en, overview_mr / overview_en).",
                 "Mobile-First Layout: Responsive grid scaling gracefully across smartphones, tablets, and 4K desktop screens.",
                 "Touch-Friendly Controls: Mobile navigation drawer, high-contrast buttons, and smooth touch scrolling."
             ])

    add_speaker_notes(slide15, "The application is built mobile-first and offers seamless bilingual switching between Marathi and English without refreshing the page.")

    # -------------------------------------------------------------
    # SLIDE 16: TESTING & VERIFICATION
    # -------------------------------------------------------------
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16, COLOR_WHITE)
    add_header(slide16, "Testing & Quality Assurance Results")

    add_card(slide16, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3),
             title="Production Build Audit",
             body_items=[
                 "npm run build executed successfully with 0 errors.",
                 "Vite 6 production bundle compiled in 18.46 seconds.",
                 "Zero broken imports or syntax errors."
             ])

    add_card(slide16, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.3),
             title="End-to-End Persistence QA",
             body_items=[
                 "100% verification across all 19 admin management modules.",
                 "Verified real-time CRUD persistence on Supabase PostgreSQL."
             ])

    add_card(slide16, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.4),
             title="Storage & Image Upload QA",
             body_items=[
                 "Direct file uploads to samarth-media bucket verified.",
                 "Automatic cache-busting versioned URLs (?v=timestamp) generated cleanly."
             ])

    add_card(slide16, Inches(6.8), Inches(4.4), Inches(5.7), Inches(2.4),
             title="Cross-Device UI Verification",
             body_items=[
                 "Mobile responsive drawer menu tested on iOS & Android.",
                 "Certificate verification search engine tested against test registration numbers."
             ])

    add_speaker_notes(slide16, "Our testing confirmed 100% data persistence across all 19 admin modules, passing production build checks with zero compilation errors.")

    # -------------------------------------------------------------
    # SLIDE 17: DEPLOYMENT
    # -------------------------------------------------------------
    slide17 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide17, COLOR_WHITE)
    add_header(slide17, "Deployment & Production Operations Architecture")

    add_card(slide17, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8),
             title="1. Code & CI/CD",
             body_items=[
                 "GitHub Repository: pavancsingh/SamarthComputer",
                 "Branch: main",
                 "Automated Cloudflare Pages build hook triggered on git push."
             ])

    add_card(slide17, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8),
             title="2. Production Build",
             body_items=[
                 "Build Command: npm run build",
                 "Output Directory: dist/",
                 "Single-Page Routing Rule: _redirects (/* /index.html 200)."
             ])

    add_card(slide17, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8),
             title="3. Live Hosting & Backend",
             body_items=[
                 "Production URL: https://samarth-computers.pages.dev",
                 "Global Edge CDN distribution.",
                 "Supabase Cloud PostgreSQL & Storage connected via environment variables."
             ])

    add_speaker_notes(slide17, "The project is deployed live on Cloudflare Pages global edge CDN, integrated directly with GitHub for CI/CD and connected securely to Supabase.")

    # -------------------------------------------------------------
    # SLIDE 18: CONCLUSION & FUTURE SCOPE
    # -------------------------------------------------------------
    slide18 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide18, COLOR_NAVY)

    txBox = slide18.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.733), Inches(1.0))
    tf = txBox.text_frame
    p_title = tf.paragraphs[0]
    p_title.text = "Conclusion & Future Roadmap"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE

    add_card(slide18, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8),
             title="Project Conclusion",
             body_items=[
                 "Modernized institute operations with a full-stack digital platform.",
                 "Eliminated paper register overhead with a real-time digital leads inbox.",
                 "Simplified citizen CSC document requirements with interactive checklists.",
                 "Delivered high-performance React 19 architecture with Supabase PostgreSQL."
             ], bg_color=RGBColor(30, 41, 59), border_color=COLOR_PRIMARY)

    add_card(slide18, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.8),
             title="Future Roadmap (v2.0)",
             body_items=[
                 "Online Fee Payment Gateway: Direct Razorpay & UPI QR code integration.",
                 "Automated SMS & WhatsApp Alerts: Instant lead receipt confirmation.",
                 "Student Learning Portal: Dedicated dashboard for typing tests & study notes.",
                 "Online Admission Workflow: Direct document upload for student admissions."
             ], bg_color=RGBColor(30, 41, 59), border_color=COLOR_EMERALD)

    add_speaker_notes(slide18, "In conclusion, Samarth Computers Khandala is a complete digital solution for computer education and government services. In future updates, we plan to add online fee payments and automated SMS alerts. Thank you!")

    prs.save(OUTPUT_PPTX)
    print(f"Presentation successfully created at: {OUTPUT_PPTX}")

if __name__ == "__main__":
    build_presentation()
